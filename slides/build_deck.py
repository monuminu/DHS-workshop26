"""Build the editable 40-slide DHS 2026 Microsoft Agent Framework deck.

The existing DHS deck is used as the visual template. Branded background artwork
is retained while all new diagrams, tables, callouts, and code are native,
editable PowerPoint objects.

Run: python slides/build_deck.py
"""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TARGET_MODE as RTM
from pptx.opc.package import _Relationship
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
OUT = HERE / "DHS-2026-Agentic-AI.pptx"
TEMPLATE = OUT

INK = RGBColor(0x18, 0x18, 0x18)
GRAY = RGBColor(0x54, 0x54, 0x5F)
MUTED = RGBColor(0x78, 0x78, 0x84)
INDIGO = RGBColor(0x4B, 0x3F, 0xC4)
MAGENTA = RGBColor(0xD6, 0x24, 0x9F)
CYAN = RGBColor(0x10, 0x9A, 0xB5)
PURPLE = RGBColor(0x6B, 0x3F, 0xA0)
GREEN = RGBColor(0x1D, 0x8B, 0x62)
RED = RGBColor(0xC9, 0x35, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xF3, 0xF1, 0xFB)
PALE_BLUE = RGBColor(0xE9, 0xF6, 0xF8)
PALE_PINK = RGBColor(0xFC, 0xEC, 0xF7)
PALE_GREEN = RGBColor(0xEA, 0xF6, 0xF0)
CODE_BG = RGBColor(0x20, 0x20, 0x2A)
FONT = "Arial"
MONO = "Aptos Mono"
SW, SH = 10.0, 5.625


def clone_slide(prs: Presentation, source_index: int):
    source = prs.slides[source_index]
    slide = prs.slides.add_slide(source.slide_layout)
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in source.shapes:
        slide.shapes._spTree.append(copy.deepcopy(shape._element))
    for r_id, rel in source.part.rels.items():
        if "image" in rel.reltype and r_id not in slide.part.rels._rels:
            slide.part.rels._rels[r_id] = _Relationship(
                slide.part.rels._base_uri,
                r_id,
                rel.reltype,
                target_mode=RTM.INTERNAL,
                target=rel.target_part,
            )
    return slide


def clean_slide(slide):
    """Keep the branded image artwork; remove prior presentation content."""
    for shape in list(slide.shapes):
        if shape.shape_type != 13:
            shape._element.getparent().remove(shape._element)


def new_slide(prs, *, dark=False):
    slide = clone_slide(prs, 0 if dark else 3)
    clean_slide(slide)
    return slide


def run_text(run, text, size=13, color=INK, *, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def textbox(slide, x, y, w, h, text="", *, size=13, color=INK, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04, font=FONT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run_text(p.add_run(), text, size, color, bold=bold, font=font)
    return shape


def title(slide, text, number, *, kicker=None):
    if kicker:
        textbox(slide, 0.55, 0.16, 8.5, 0.22, kicker.upper(), size=8.5,
                color=MAGENTA, bold=True)
    textbox(slide, 0.55, 0.38, 8.65, 0.58, text, size=23, bold=True,
            valign=MSO_ANCHOR.MIDDLE)
    textbox(slide, 9.22, 0.45, 0.3, 0.3, f"{number:02d}", size=9, color=MUTED,
            bold=True, align=PP_ALIGN.RIGHT)


def box(slide, x, y, w, h, *, fill=PALE, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def label(slide, x, y, w, text, *, fill=INDIGO, color=WHITE, size=10):
    shape = box(slide, x, y, w, 0.34, fill=fill)
    shape.text_frame.clear()
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_text(p.add_run(), text, size, color, bold=True)
    return shape


def bullets(slide, x, y, w, h, items, *, size=12.5, color=INK, gap=4):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = level
        glyph = "–  " if level else "•  "
        run_text(p.add_run(), glyph, size, INDIGO if level == 0 else MAGENTA, bold=True)
        run_text(p.add_run(), text, size - (0.8 if level else 0), color if level == 0 else GRAY)
    return shape


def card(slide, x, y, w, h, heading, body, *, accent=INDIGO, fill=PALE,
         heading_size=12, body_size=10.5):
    box(slide, x, y, w, h, fill=fill)
    box(slide, x, y, 0.08, h, fill=accent, radius=False)
    textbox(slide, x + 0.18, y + 0.12, w - 0.3, 0.35, heading, size=heading_size,
            bold=True, color=accent)
    textbox(slide, x + 0.18, y + 0.52, w - 0.3, h - 0.62, body, size=body_size,
            color=GRAY)


def code(slide, x, y, w, h, text, *, caption=None, size=9.2):
    box(slide, x, y, w, h, fill=CODE_BG)
    if caption:
        textbox(slide, x + 0.16, y + 0.08, w - 0.32, 0.25, caption, size=8,
                color=CYAN, bold=True, font=MONO)
        offset = 0.34
    else:
        offset = 0.13
    textbox(slide, x + 0.14, y + offset, w - 0.28, h - offset - 0.1, text,
            size=size, color=WHITE, font=MONO)


def arrow(slide, x1, y1, x2, y2, *, color=INDIGO, width=2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def flow(slide, y, steps, *, x=0.65, total_w=8.7, colors=None, height=0.68):
    colors = colors or [INDIGO, MAGENTA, CYAN, PURPLE, GREEN]
    gap = 0.28
    width = (total_w - gap * (len(steps) - 1)) / len(steps)
    for i, step in enumerate(steps):
        xx = x + i * (width + gap)
        shape = box(slide, xx, y, width, height, fill=colors[i % len(colors)])
        shape.text_frame.clear()
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run_text(p.add_run(), step, 10.5, WHITE, bold=True)
        if i < len(steps) - 1:
            arrow(slide, xx + width + 0.02, y + height / 2,
                  xx + width + gap - 0.02, y + height / 2, color=GRAY, width=1.5)


def table(slide, x, y, w, h, headers, rows, widths=None, *, font_size=9.5):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    tbl = shape.table
    if widths:
        for i, value in enumerate(widths):
            tbl.columns[i].width = Inches(value)
    for r in range(len(rows) + 1):
        for c in range(len(headers)):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = INDIGO if r == 0 else (PALE if r % 2 else WHITE)
            cell.text = headers[c] if r == 0 else str(rows[r - 1][c])
            for p in cell.text_frame.paragraphs:
                for rn in p.runs:
                    rn.font.name = FONT
                    rn.font.size = Pt(font_size if r else font_size - 0.3)
                    rn.font.bold = r == 0
                    rn.font.color.rgb = WHITE if r == 0 else INK
    return shape


def callout(slide, x, y, w, text, *, accent=MAGENTA, fill=PALE_PINK, size=10.5):
    box(slide, x, y, w, 0.55, fill=fill, line=accent)
    textbox(slide, x + 0.12, y + 0.09, w - 0.24, 0.34, text, size=size,
            color=accent, bold=True, valign=MSO_ANCHOR.MIDDLE)


def actor_sequence(slide, actors, events, *, y=1.45):
    x0, span = 0.75, 8.5
    xs = [x0 + i * span / (len(actors) - 1) for i in range(len(actors))]
    for xx, actor in zip(xs, actors):
        label(slide, xx - 0.55, y, 1.1, actor, fill=INDIGO, size=9)
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(xx), Inches(y + 0.4),
            Inches(xx), Inches(y + 3.4),
        )
        line.line.color.rgb = RGBColor(0xC8, 0xC8, 0xD0)
        line.line.dash_style = 2
    yy = y + 0.62
    for source, target, event, color in events:
        x1, x2 = xs[source], xs[target]
        arrow(slide, x1, yy, x2, yy, color=color, width=1.6)
        textbox(slide, min(x1, x2) + 0.1, yy - 0.25, abs(x2 - x1) - 0.2, 0.22,
                event, size=8.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        yy += 0.48


def build():
    prs = Presentation(TEMPLATE)
    template_count = len(prs.slides)
    built = []

    # 1
    s = new_slide(prs, dark=True); built.append(s)
    textbox(s, 0.72, 1.05, 8.45, 1.45,
            "Building Agents with the\nMicrosoft Agent Framework",
            size=27, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    textbox(s, 0.75, 2.72, 7.6, 0.66,
            "A deep technical workshop on loops, context, harnesses,\norchestration, evaluation, and production operation",
            size=14, color=RGBColor(0xE4, 0xE4, 0xEA))
    label(s, 0.75, 3.72, 1.25, "DHS 2026", fill=MAGENTA)
    textbox(s, 2.15, 3.73, 5.7, 0.35, "Manoranjan Rajguru  ·  Microsoft", size=11,
            color=WHITE, bold=True)
    textbox(s, 0.75, 4.75, 8.5, 0.3,
            "MODEL → TOOLS → STATE → CONTROL FLOW → EVIDENCE → OPERATIONS",
            size=10, color=CYAN, bold=True)

    # 2
    s = new_slide(prs); built.append(s); title(s, "Why Agentic Systems Now", 2, kicker="Foundations")
    card(s, 0.55, 1.22, 2.8, 1.35, "Capability moved", "Models can select tools, follow structured schemas, and reason across long task horizons.", accent=INDIGO)
    card(s, 3.6, 1.22, 2.8, 1.35, "Complexity moved outward", "Reliability now depends on the execution layer: state, retries, approvals, persistence, and budgets.", accent=MAGENTA, fill=PALE_PINK)
    card(s, 6.65, 1.22, 2.8, 1.35, "Value moved to systems", "Production differentiation is the harness around the model—not another prompt wrapper.", accent=CYAN, fill=PALE_BLUE)
    flow(s, 3.18, ["PROMPT", "TOOL USE", "DURABLE LOOP", "GOVERNED SYSTEM"])
    callout(s, 0.7, 4.35, 8.6, "Engineering question: what must remain deterministic when the model is not?", accent=PURPLE)

    # 3
    s = new_slide(prs); built.append(s); title(s, "Workshop Goals and Expected Outcomes", 3, kicker="System outcomes")
    table(s, 0.6, 1.2, 8.8, 3.55,
          ["By the end, you can…", "Artifact", "Engineering decision"],
          [
              ("Build", "Agent + typed tools", "Define the loop and provider boundary"),
              ("Engineer context", "Session + providers", "Choose injection and persistence boundaries"),
              ("Assemble", "Harness agent", "Enable only the batteries the workload needs"),
              ("Coordinate", "Workflow / orchestration", "Prefer explicit control when recovery matters"),
              ("Measure", "Golden dataset + evaluators", "Separate CI gates from quality signals"),
              ("Operate", "OTel + middleware", "Trace every model/tool/state transition"),
          ], widths=[3.1, 2.2, 3.5], font_size=9.5)
    callout(s, 0.75, 4.92, 8.5, "Target: a recoverable, observable, provider-portable agent—not a one-shot demo.", accent=INDIGO, fill=PALE)

    # 4
    s = new_slide(prs); built.append(s); title(s, "LLM Call vs Agent", 4, kicker="Mental model")
    card(s, 0.6, 1.25, 3.65, 2.85, "LLM CALL", "Input → model → output\n\nStateless by default\nNo action semantics\nSingle failure boundary\nCaller owns retries and history", accent=MUTED, fill=WHITE, body_size=12)
    card(s, 5.75, 1.25, 3.65, 2.85, "AGENT", "Conversation + tools → repeated model calls\n\nActs through typed interfaces\nMutates session state\nStops on a completion condition\nHarness owns lifecycle policy", accent=INDIGO, body_size=12)
    arrow(s, 4.45, 2.6, 5.55, 2.6, color=MAGENTA, width=3)
    textbox(s, 4.4, 2.07, 1.2, 0.35, "ADD A LOOP", size=9, color=MAGENTA, bold=True, align=PP_ALIGN.CENTER)
    callout(s, 0.75, 4.55, 8.5, "An agent.run() may make N model calls and M tool calls before returning once.", accent=PURPLE)

    # 5
    s = new_slide(prs); built.append(s); title(s, "Agent Anatomy", 5, kicker="Definition")
    flow(s, 1.38, ["MODEL", "INSTRUCTIONS", "TOOLS", "LOOP"], x=0.65, total_w=8.7, height=0.78)
    textbox(s, 0.75, 2.48, 8.5, 0.55,
            "Agent  =  model  +  instructions  +  tools  +  a loop that runs until the task is done",
            size=18, color=INDIGO, bold=True, align=PP_ALIGN.CENTER)
    card(s, 0.7, 3.35, 2.0, 1.15, "MODEL", "Probabilistic planner and generator", accent=INDIGO)
    card(s, 2.9, 3.35, 2.0, 1.15, "POLICY", "Instructions and approval boundaries", accent=MAGENTA, fill=PALE_PINK)
    card(s, 5.1, 3.35, 2.0, 1.15, "EFFECTORS", "Typed functions and external services", accent=CYAN, fill=PALE_BLUE)
    card(s, 7.3, 3.35, 2.0, 1.15, "RUNTIME", "State, termination, persistence, telemetry", accent=PURPLE)

    # 6
    s = new_slide(prs); built.append(s); title(s, "The Agent Loop / ReAct Cycle", 6, kicker="Lifecycle")
    actor_sequence(s, ["HOST", "AGENT", "MODEL", "TOOL"], [
        (0, 1, "run(input, session)", INDIGO),
        (1, 2, "messages + tool schemas", MAGENTA),
        (2, 1, "function call", CYAN),
        (1, 3, "validate + invoke", PURPLE),
        (3, 1, "function result", GREEN),
        (1, 2, "append observation", MAGENTA),
        (2, 1, "final answer / next call", CYAN),
    ])
    callout(s, 0.75, 4.95, 8.5, "Termination is observable: no further tool call, explicit policy stop, budget exhaustion, or failure.", accent=RED, fill=PALE_PINK, size=9.8)

    # 7
    s = new_slide(prs); built.append(s); title(s, "Why Tools Matter", 7, kicker="Action boundary")
    bullets(s, 0.6, 1.25, 4.3, 3.55, [
        "Tools turn language into typed, auditable effects.",
        "Schemas constrain names, arguments, and result shape.",
        "Descriptions are routing policy: ambiguity creates bad calls.",
        "Small, single-purpose tools compose better than “do_everything”.",
        "Side effects require idempotency, authorization, and approval.",
    ], size=12.4, gap=7)
    card(s, 5.2, 1.25, 4.15, 0.9, "READ", "Search, retrieve, inspect—usually auto-approve.", accent=GREEN, fill=PALE_GREEN)
    card(s, 5.2, 2.35, 4.15, 0.9, "COMPUTE", "Transform, rank, calculate—bound CPU/time.", accent=CYAN, fill=PALE_BLUE)
    card(s, 5.2, 3.45, 4.15, 0.9, "WRITE", "Send, delete, deploy—require policy + human gate.", accent=RED, fill=PALE_PINK)
    callout(s, 5.25, 4.62, 4.05, "Treat tool output as untrusted input.", accent=PURPLE, size=9.6)

    # 8
    s = new_slide(prs); built.append(s); title(s, "Why Memory and Context Matter", 8, kicker="Information state")
    flow(s, 1.28, ["SYSTEM POLICY", "SESSION HISTORY", "MEMORY", "SKILLS", "TOOL RESULTS"])
    textbox(s, 0.7, 2.28, 8.6, 0.5, "The model does not “remember”; the runtime reconstructs a context window for every call.", size=16, color=INDIGO, bold=True, align=PP_ALIGN.CENTER)
    card(s, 0.7, 3.1, 2.55, 1.35, "Too little", "Repeated questions, lost commitments, poor continuity.", accent=RED, fill=PALE_PINK)
    card(s, 3.7, 3.1, 2.55, 1.35, "Too much", "Token waste, attention dilution, stale or contradictory state.", accent=MAGENTA)
    card(s, 6.7, 3.1, 2.55, 1.35, "Engineered", "Relevant, ordered, attributable, budgeted information.", accent=GREEN, fill=PALE_GREEN)

    # 9
    s = new_slide(prs); built.append(s); title(s, "Context Engineering Levers", 9, kicker="M3")
    table(s, 0.6, 1.2, 8.8, 3.65,
          ["Lever", "Framework surface", "Boundary / tradeoff"],
          [
              ("Sessions", "agent.create_session()", "Continuity; grows unless compacted"),
              ("Providers", "ContextProvider", "Dynamic injection + state mutation"),
              ("Memory", "Provider state / MemoryStore", "Selective durability, privacy lifecycle"),
              ("Skills", "SkillsProvider", "Cheap catalog; load details on demand"),
              ("Compaction", "CompactionProvider", "Fit budget; lossy by definition"),
          ], widths=[1.45, 2.55, 4.8], font_size=10.2)
    callout(s, 0.75, 4.98, 8.5, "Order matters: policy → durable facts → working state → recent conversation → observations.", accent=INDIGO)

    # 10
    s = new_slide(prs); built.append(s); title(s, "Skills and Progressive Disclosure", 10, kicker="Context economics")
    card(s, 0.65, 1.2, 2.2, 3.35, "ALWAYS IN CONTEXT", "Skill name\nOne-line description\nTrigger hints\n\nCost scales with catalog entries, not reference size.", accent=INDIGO, body_size=11.5)
    card(s, 3.4, 1.2, 2.2, 3.35, "ON DEMAND", "SKILL.md instructions\nReference documents\nScripts\n\nLoaded only after the model selects the skill.", accent=MAGENTA, fill=PALE_PINK, body_size=11.5)
    card(s, 6.15, 1.2, 3.2, 3.35, "DESIGN CONSEQUENCE", "Fifty skills become fifty short index lines—not fifty pages.\n\nApply the same pattern to memory: inject an index first; fetch the full fact only when relevant.", accent=CYAN, fill=PALE_BLUE, body_size=11.2)
    arrow(s, 2.88, 2.7, 3.32, 2.7, color=MAGENTA)
    arrow(s, 5.63, 2.7, 6.07, 2.7, color=CYAN)

    # 11
    s = new_slide(prs); built.append(s); title(s, "From Primitives to Harness", 11, kicker="Composition")
    flow(s, 1.25, ["AGENT LOOP", "+ HISTORY", "+ COMPACTION", "+ PLANNING", "+ TELEMETRY"])
    bullets(s, 0.65, 2.35, 4.25, 2.55, [
        "Hand-assembly is useful for learning and narrow services.",
        "Production agents repeatedly need the same lifecycle machinery.",
        "A harness standardizes defaults without hiding extension points.",
    ], size=12.5, gap=9)
    card(s, 5.2, 2.25, 4.1, 2.2, "HARNESS = EXECUTION POLICY", "What runs before/after each call\nWhere state is authoritative\nWhen work persists\nHow context is budgeted\nWhich operations require approval", accent=PURPLE, body_size=11.2)
    callout(s, 0.75, 4.82, 8.5, "Abstraction exists to make the non-model behavior explicit and repeatable.", accent=INDIGO)

    # 12
    s = new_slide(prs); built.append(s); title(s, "Microsoft Agent Framework Architecture Overview", 12, kicker="Framework map")
    layers = [
        ("APPLICATION / HOST", "API · notebook · A2A · Azure Functions · container", PURPLE),
        ("AGENTS + ORCHESTRATIONS", "Agent · harness agent · sequential · concurrent · handoff · group chat", INDIGO),
        ("WORKFLOW RUNTIME", "WorkflowBuilder · typed executors · edges · checkpoints · recovery", MAGENTA),
        ("RUNTIME EXTENSIONS", "Tools · context providers · middleware · approvals · OpenTelemetry", CYAN),
        ("MODEL PROVIDERS", "Foundry · OpenAI · Azure OpenAI · Anthropic · Ollama · Bedrock · Gemini", GREEN),
    ]
    y = 1.12
    for name, detail, accent in layers:
        box(s, 0.85, y, 8.3, 0.66, fill=WHITE, line=accent)
        label(s, 1.0, y + 0.15, 2.15, name, fill=accent, size=8.5)
        textbox(s, 3.35, y + 0.17, 5.45, 0.27, detail, size=9.5, color=GRAY)
        y += 0.78

    # 13
    s = new_slide(prs); built.append(s); title(s, "Provider-Agnostic Model Switching", 13, kicker="M1 · Boundary")
    code(s, 0.6, 1.18, 4.55, 3.55, '''def get_chat_client():
    provider = os.getenv("MODEL_PROVIDER", "foundry")
    if provider == "foundry":
        return FoundryChatClient(...)
    if provider == "openai":
        return OpenAIChatClient(...)
    if provider == "azure-openai":
        return AzureOpenAIChatClient(...)
    raise ValueError(f"Unsupported: {provider}")''', caption="workshop_utils/clients.py", size=9.0)
    card(s, 5.45, 1.18, 3.95, 1.05, "ONE SWITCH", "MODEL_PROVIDER changes deployment—not agent code.", accent=INDIGO)
    card(s, 5.45, 2.43, 3.95, 1.05, "STABLE CONTRACT", "Agent consumes ChatClient; provider auth/options stay at the edge.", accent=CYAN, fill=PALE_BLUE)
    card(s, 5.45, 3.68, 3.95, 1.05, "TEST THE SEAMS", "Structured output, tool schemas, streaming, and usage metadata vary.", accent=MAGENTA, fill=PALE_PINK)

    # 14
    s = new_slide(prs); built.append(s); title(s, "First Agent: Minimal Code Walkthrough", 14, kicker="M1 · Construction")
    code(s, 0.6, 1.15, 5.1, 3.7, '''from agent_framework import Agent
from workshop_utils.clients import get_chat_client

agent = Agent(
    client=get_chat_client(),
    name="architecture-reviewer",
    instructions=(
        "Review designs for failure boundaries, "
        "state ownership, and observability."
    ),
)

response = await agent.run("Review this agent topology.")
print(response.text)''', caption="minimal agent", size=9.5)
    card(s, 5.95, 1.15, 3.45, 1.0, "client", "The provider boundary and model transport.", accent=INDIGO)
    card(s, 5.95, 2.35, 3.45, 1.0, "instructions", "Durable behavior policy—not task input.", accent=MAGENTA, fill=PALE_PINK)
    card(s, 5.95, 3.55, 3.45, 1.0, "run()", "One host turn; internally N service calls.", accent=CYAN, fill=PALE_BLUE)
    callout(s, 5.98, 4.75, 3.38, "Name agents for traces and delegation.", accent=PURPLE, size=9.2)

    # 15
    s = new_slide(prs); built.append(s); title(s, "Streaming vs Non-Streaming Execution", 15, kicker="M1 · Response semantics")
    code(s, 0.6, 1.2, 4.25, 1.45, '''result = await agent.run(prompt)
print(result.text)''', caption="buffered")
    code(s, 0.6, 3.0, 4.25, 1.65, '''async for update in agent.run(prompt, stream=True):
    if update.text:
        print(update.text, end="", flush=True)''', caption="streaming", size=8.8)
    table(s, 5.15, 1.2, 4.25, 3.45,
          ["Concern", "Buffered", "Streaming"],
          [
              ("First token", "After completion", "As produced"),
              ("UX", "Batch / jobs", "Interactive"),
              ("Errors", "Single terminal path", "Partial output possible"),
              ("Persistence", "End-of-turn", "Still per service call"),
              ("Consumers", "Simple result", "Backpressure + cancellation"),
          ], widths=[1.45, 1.4, 1.4], font_size=8.9)
    callout(s, 0.75, 4.92, 8.5, "Streaming changes delivery semantics; it does not bypass the agent loop.", accent=INDIGO)

    # 16
    s = new_slide(prs); built.append(s); title(s, "Function Tools: Anatomy of a Tool", 16, kicker="M2 · Typed effects")
    code(s, 0.6, 1.18, 4.9, 3.7, '''from agent_framework import tool

@tool(approval_mode="never_require")
def get_order(order_id: str) -> dict:
    """Return status and ETA for one order.

    Args:
        order_id: Canonical order identifier.
    """
    return order_store.lookup(order_id)

agent = Agent(client=client, tools=[get_order])''', caption="schema generated from signature + docstring", size=9.0)
    card(s, 5.8, 1.18, 3.55, 0.95, "NAME", "Routing token; make it precise and stable.", accent=INDIGO)
    card(s, 5.8, 2.3, 3.55, 0.95, "SCHEMA", "Types constrain arguments; validate again at runtime.", accent=CYAN, fill=PALE_BLUE)
    card(s, 5.8, 3.42, 3.55, 0.95, "DESCRIPTION", "Tells the model when—and when not—to call.", accent=MAGENTA, fill=PALE_PINK)
    callout(s, 5.83, 4.55, 3.49, "Return machine-readable results.", accent=GREEN, fill=PALE_GREEN, size=9.4)

    # 17
    s = new_slide(prs); built.append(s); title(s, "Tool-Calling Lifecycle and Approvals", 17, kicker="M2 · Control point")
    actor_sequence(s, ["MODEL", "RUNTIME", "POLICY", "TOOL"], [
        (0, 1, "call(name, args)", INDIGO),
        (1, 2, "authorize + approval?", MAGENTA),
        (2, 1, "approve / deny / edit", RED),
        (1, 3, "validated invocation", CYAN),
        (3, 1, "result / exception", GREEN),
        (1, 0, "function_result", PURPLE),
    ], y=1.35)
    label(s, 0.75, 4.63, 2.55, 'always_require', fill=RED)
    label(s, 3.72, 4.63, 2.55, 'never_require', fill=GREEN)
    label(s, 6.68, 4.63, 2.55, 'policy-based', fill=PURPLE)
    textbox(s, 0.85, 5.02, 8.3, 0.25, "High-impact writes  ·  safe deterministic reads  ·  runtime rule by identity, arguments, environment", size=8.8, color=GRAY, align=PP_ALIGN.CENTER)

    # 18
    s = new_slide(prs); built.append(s); title(s, "Multi-Tool Composition Patterns", 18, kicker="M2 · Planning")
    table(s, 0.6, 1.18, 8.8, 3.6,
          ["Pattern", "Shape", "Use when", "Primary failure mode"],
          [
              ("Serial dependency", "A → B → C", "Output feeds next call", "Cascading bad argument"),
              ("Parallel reads", "A ∥ B ∥ C", "Independent evidence", "Rate limit / partial failure"),
              ("Read–decide–write", "R → policy → W", "Side effect follows evidence", "Stale read / race"),
              ("Compensating action", "W → verify → undo", "No distributed transaction", "Compensation also fails"),
          ], widths=[1.55, 1.35, 2.75, 3.15], font_size=9.4)
    callout(s, 0.75, 4.98, 8.5, "Do not encode critical ordering only in prose: use a workflow when sequence is a correctness property.", accent=RED, fill=PALE_PINK, size=9.8)

    # 19
    s = new_slide(prs); built.append(s); title(s, "MCP / External Tool Integration Considerations", 19, kicker="M2 · Protocol boundary")
    card(s, 0.6, 1.2, 2.7, 1.45, "DISCOVERY", "Server exposes tool catalog and schemas. Cache carefully; capabilities can change.", accent=INDIGO)
    card(s, 3.65, 1.2, 2.7, 1.45, "TRANSPORT", "stdio, HTTP, or WebSocket. Bound startup, request, and idle timeouts.", accent=CYAN, fill=PALE_BLUE)
    card(s, 6.7, 1.2, 2.7, 1.45, "TRUST", "Remote descriptions and results are untrusted. Apply allowlists and output limits.", accent=RED, fill=PALE_PINK)
    card(s, 0.6, 3.0, 2.7, 1.45, "IDENTITY", "Propagate least-privilege credentials; do not hand the model raw secrets.", accent=PURPLE)
    card(s, 3.65, 3.0, 2.7, 1.45, "VERSIONING", "Pin protocol/package compatibility and test schema drift.", accent=MAGENTA, fill=PALE_PINK)
    card(s, 6.7, 3.0, 2.7, 1.45, "OBSERVABILITY", "Trace server, tool, arguments metadata, latency, and sanitized result size.", accent=GREEN, fill=PALE_GREEN)
    callout(s, 0.75, 4.85, 8.5, "MCP standardizes connectivity—not authorization, safety, or reliability.", accent=INDIGO)

    # 20
    s = new_slide(prs); built.append(s); title(s, "Sessions and Persistent Conversational State", 20, kicker="M3 · State boundary")
    code(s, 0.6, 1.18, 4.35, 2.15, '''session = agent.create_session()

await agent.run("My region is eu-west.", session=session)
await agent.run("Use the same region.", session=session)

# New session: no conversational continuity
fresh = agent.create_session()''', caption="session reuse")
    card(s, 5.25, 1.18, 4.1, 1.05, "SESSION HISTORY", "Turn-by-turn messages; useful but unbounded and ephemeral unless persisted.", accent=INDIGO)
    card(s, 5.25, 2.43, 4.1, 1.05, "SESSION STATE", "Provider-owned structured state: todos, mode, memory indexes.", accent=CYAN, fill=PALE_BLUE)
    card(s, 5.25, 3.68, 4.1, 1.05, "DURABLE STORE", "External lifecycle: retention, encryption, tenant isolation, restore.", accent=PURPLE)
    callout(s, 0.65, 4.35, 4.25, "Agent object = behavior; session = one continuity boundary.", accent=MAGENTA, size=9.4)

    # 21
    s = new_slide(prs); built.append(s); title(s, "ContextProvider Deep Dive", 21, kicker="M3 · Injection + mutation")
    code(s, 0.55, 1.12, 5.4, 3.95, '''class UserMemoryProvider(ContextProvider):
    DEFAULT_SOURCE_ID = "user_memory"

    async def before_run(self, *, context, state, **_):
        if name := state.get("user_name"):
            context.extend_instructions(
                self.source_id, f"User name: {name}"
            )

    async def after_run(self, *, context, state, **_):
        for msg in context.input_messages:
            if name := extract_name(msg.text):
                state["user_name"] = name''', caption="two hooks around the loop", size=8.65)
    card(s, 6.2, 1.12, 3.2, 1.45, "before_run", "Read provider state and inject attributable instructions/messages into the next call.", accent=INDIGO)
    card(s, 6.2, 2.82, 3.2, 1.45, "after_run", "Observe the exchange and mutate selective state after the turn.", accent=MAGENTA, fill=PALE_PINK)
    callout(s, 6.23, 4.55, 3.14, "Use source_id for provenance.", accent=CYAN, fill=PALE_BLUE, size=9.2)

    # 22
    s = new_slide(prs); built.append(s); title(s, "Durable Memory Design Patterns", 22, kicker="M3 · Persistence")
    table(s, 0.6, 1.18, 8.8, 3.55,
          ["Pattern", "Persist", "Inject", "Risk / mitigation"],
          [
              ("Profile facts", "Stable preferences", "Small verified subset", "Stale facts → timestamp + edit"),
              ("Episodic summary", "Outcome + decisions", "Relevant episode only", "Lossy → retain source pointer"),
              ("Semantic index", "Chunks + embeddings", "Top-k with citations", "Poisoning → trust metadata"),
              ("File memory", "Human-readable files", "Index then fetch", "Tenant leakage → scoped roots"),
          ], widths=[1.55, 2.1, 2.2, 2.95], font_size=9.4)
    callout(s, 0.75, 4.92, 8.5, "Progressive disclosure: expose a compact memory index; load the full record only when selected.", accent=INDIGO)

    # 23
    s = new_slide(prs); built.append(s); title(s, "Compaction and Context-Window Management", 23, kicker="M3 · Token budget")
    flow(s, 1.25, ["MEASURE", "PRIORITIZE", "SUMMARIZE / TRIM", "PERSIST", "REHYDRATE"])
    card(s, 0.65, 2.35, 2.65, 2.15, "BEFORE EACH CALL", "Fit outgoing messages to max_context_window_tokens.\n\nPreserve policy, active todos, recent tool evidence, and unresolved commitments.", accent=INDIGO, body_size=10.8)
    card(s, 3.68, 2.35, 2.65, 2.15, "AFTER EACH TURN", "Compact persisted history in place so future calls start from a bounded representation.", accent=MAGENTA, fill=PALE_PINK, body_size=10.8)
    card(s, 6.7, 2.35, 2.65, 2.15, "FAILURE MODES", "Summary drops IDs or constraints; stale summary outranks fresh evidence; no token budget means compaction is off.", accent=RED, fill=PALE_PINK, body_size=10.4)
    callout(s, 0.75, 4.82, 8.5, "Compaction is lossy compression—evaluate task continuity, not just token count.", accent=PURPLE)

    # 24
    s = new_slide(prs); built.append(s); title(s, "Harness Overview: create_harness_agent", 24, kicker="M4 · Batteries included")
    code(s, 0.55, 1.15, 4.7, 3.85, '''agent = create_harness_agent(
    client=get_chat_client(),
    name="researcher",
    instructions=POLICY,
    default_options={"store": False},
    max_context_window_tokens=64_000,
    max_tool_result_tokens=8_000,
    memory_path=MEMORY,
    skills_path=SKILLS,
    enable_background_agents=True,
    autonomous=True,
)''', caption="one factory; explicit policy", size=8.65)
    table(s, 5.52, 1.15, 3.9, 3.85,
          ["Battery", "Purpose"],
          [
              ("Function invocation", "Tool loop"),
              ("History", "Local continuity"),
              ("Compaction", "Token budget"),
              ("TodoProvider", "Plan state"),
              ("AgentModeProvider", "Plan / execute"),
              ("File memory + access", "Durable files"),
              ("SkillsProvider", "On-demand expertise"),
              ("OpenTelemetry", "Trace lifecycle"),
          ], widths=[1.72, 2.18], font_size=8.3)

    # 25
    s = new_slide(prs); built.append(s); title(s, "Local-Authoritative History and Crash Safety", 25, kicker="M4 · Durability contract")
    actor_sequence(s, ["HOST", "HARNESS", "MODEL", "LOCAL STORE"], [
        (0, 1, "run(turn)", INDIGO),
        (1, 2, "call #1  store=False", MAGENTA),
        (2, 1, "tool request", CYAN),
        (1, 3, "persist immediately", GREEN),
        (1, 2, "call #2 + tool result", MAGENTA),
        (2, 1, "final response", CYAN),
        (1, 3, "persist immediately", GREEN),
    ], y=1.3)
    callout(s, 0.65, 4.83, 4.15, "Local history is authoritative for compaction, restore, and provider portability.", accent=INDIGO, size=9.2)
    callout(s, 5.05, 4.83, 4.3, "require_per_service_call_history_persistence=True limits crash loss to the in-flight call.", accent=GREEN, fill=PALE_GREEN, size=9.0)

    # 26
    s = new_slide(prs); built.append(s); title(s, "TodoProvider and the Planning Model", 26, kicker="M4 · Explicit work state")
    card(s, 0.6, 1.2, 2.6, 3.55, "TOOLS INJECTED", "todos_add\ntodos_complete\ntodos_remove\ntodos_get_remaining\ntodos_get_all", accent=INDIGO, body_size=12)
    card(s, 3.7, 1.2, 2.6, 3.55, "STATE MODEL", "Todo state belongs to the session.\n\nThe current list is re-injected before every model call, so plan state survives attention shifts.", accent=MAGENTA, fill=PALE_PINK, body_size=11.2)
    card(s, 6.8, 1.2, 2.6, 3.55, "HOST CONTROL", "Render progress externally.\nResume unfinished work.\nEnforce limits.\nRequire evidence before completion.\nStop when no items remain.", accent=CYAN, fill=PALE_BLUE, body_size=11.0)
    callout(s, 0.75, 4.92, 8.5, "A todo list is observable control state—not proof that the work is correct.", accent=PURPLE)

    # 27
    s = new_slide(prs); built.append(s); title(s, "Plan vs Execute Mode", 27, kicker="M4 · State machine")
    box(s, 0.8, 1.35, 3.3, 2.45, fill=PALE_BLUE, line=CYAN)
    textbox(s, 1.0, 1.6, 2.9, 0.45, "PLAN", size=21, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 1.05, 2.25, 2.8, 1.0, "Explore constraints\nCreate / revise todos\nNo irreversible side effects", size=12, color=GRAY, align=PP_ALIGN.CENTER)
    box(s, 5.9, 1.35, 3.3, 2.45, fill=PALE_GREEN, line=GREEN)
    textbox(s, 6.1, 1.6, 2.9, 0.45, "EXECUTE", size=21, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 6.15, 2.25, 2.8, 1.0, "Perform approved work\nComplete todos with evidence\nAutonomous loop may continue", size=12, color=GRAY, align=PP_ALIGN.CENTER)
    arrow(s, 4.15, 2.05, 5.82, 2.05, color=MAGENTA, width=3)
    arrow(s, 5.82, 3.15, 4.15, 3.15, color=PURPLE, width=2)
    textbox(s, 4.35, 1.55, 1.3, 0.32, "APPROVE", size=9, color=MAGENTA, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 4.3, 3.3, 1.4, 0.32, "REPLAN / PAUSE", size=8.5, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    callout(s, 0.85, 4.35, 8.3, "External mode changes inject a notice on the next turn—history alone may anchor the model to its previous mode.", accent=INDIGO, size=9.7)

    # 28
    s = new_slide(prs); built.append(s); title(s, "File Memory and File Access Providers", 28, kicker="M4 · Durable workspace")
    card(s, 0.6, 1.2, 4.1, 2.9, "FileMemoryProvider", "Durable memory across sessions\nHuman-inspectable and versionable\nSurvives conversation compaction\nInject index first, content on demand\n\nUse for curated facts—not raw transcripts.", accent=INDIGO, body_size=11.3)
    card(s, 5.3, 1.2, 4.1, 2.9, "FileAccessProvider", "Scoped read/write operations\nWorkspace as an agent artifact boundary\nHost can inspect outputs independently\n\nCanonicalize paths; deny traversal; separate tenant roots; cap file size.", accent=CYAN, fill=PALE_BLUE, body_size=11.3)
    callout(s, 0.75, 4.45, 8.5, "Filesystem access is a capability boundary: root scope, extension policy, quotas, and audit logs are mandatory.", accent=RED, fill=PALE_PINK, size=9.5)

    # 29
    s = new_slide(prs); built.append(s); title(s, "SkillsProvider and On-Demand Loading", 29, kicker="M4 · Progressive expertise")
    code(s, 0.55, 1.15, 4.75, 3.8, '''skills = SkillsProvider(
    FileSkillsSource("skills/")
)

agent = create_harness_agent(
    client=get_chat_client(),
    skills_provider=skills,
    auto_approval_rules=[
        SkillsProvider.all_tools_auto_approval_rule
    ],
)''', caption="catalog → select → load → use", size=8.9)
    flow(s, 1.35, ["DISCOVER", "SELECT", "LOAD", "EXECUTE"], x=5.65, total_w=3.65, height=0.55)
    card(s, 5.55, 2.25, 3.85, 1.0, "CATALOG", "Names + descriptions stay in the prompt.", accent=INDIGO)
    card(s, 5.55, 3.45, 3.85, 1.0, "DETAIL", "Instructions, references, and scripts load only when chosen.", accent=MAGENTA, fill=PALE_PINK)
    callout(s, 5.58, 4.65, 3.79, "Auto-approve loading ≠ auto-approve skill side effects.", accent=RED, fill=PALE_PINK, size=8.7)

    # 30
    s = new_slide(prs); built.append(s); title(s, "Background Agents and Autonomous Looping", 30, kicker="M4 · Long-running work")
    flow(s, 1.18, ["PLAN TODOS", "DELEGATE", "COLLECT", "COMPLETE", "CONTINUE?"])
    code(s, 0.6, 2.25, 4.3, 2.35, '''agent = create_harness_agent(
    ...,
    enable_background_agents=True,
    autonomous=True,
    loop_should_continue=todos_remaining(
        looping_modes=["execute"]
    ),
)''', caption="explicit continuation predicate", size=8.8)
    card(s, 5.2, 2.25, 4.2, 2.35, "BOUND THE LOOP", "Max iterations / wall time\nToken + tool budgets\nCancellation propagation\nIdempotent resumability\nHeartbeat + lease ownership\nEscalate repeated failures", accent=RED, fill=PALE_PINK, body_size=10.9)
    callout(s, 0.75, 4.85, 8.5, "One run() is one turn; autonomous completion requires a host-visible continuation and stop policy.", accent=INDIGO)

    # 31
    s = new_slide(prs); built.append(s); title(s, "OpenTelemetry and Observability Architecture", 31, kicker="M7 · Operations")
    layers = [
        ("SCENARIO SPAN", "user request · tenant · session · release", PURPLE),
        ("AGENT / WORKFLOW", "mode · todo count · checkpoint · delegation", INDIGO),
        ("MODEL CALL", "provider · model · latency · tokens · finish reason", MAGENTA),
        ("TOOL CALL", "tool · approval · latency · status · result size", CYAN),
        ("STATE / COMPACTION", "persist · restore · token budget · summary", GREEN),
    ]
    y = 1.12
    for name, details, accent in layers:
        box(s, 0.65, y, 5.55, 0.64, fill=WHITE, line=accent)
        label(s, 0.82, y + 0.14, 1.72, name, fill=accent, size=8)
        textbox(s, 2.75, y + 0.16, 3.25, 0.28, details, size=8.8, color=GRAY)
        y += 0.76
    card(s, 6.55, 1.12, 2.8, 1.25, "MIDDLEWARE", "Control point before/after each model call: usage, redaction, retries, policy.", accent=INDIGO)
    card(s, 6.55, 2.65, 2.8, 1.25, "OTEL", "Portable traces emitted to console, Phoenix, Langfuse, or another OTLP backend.", accent=CYAN, fill=PALE_BLUE)
    card(s, 6.55, 4.18, 2.8, 0.78, "DATA POLICY", "Sensitive content off by default.", accent=RED, fill=PALE_PINK, body_size=9.3)

    # 32
    s = new_slide(prs); built.append(s); title(s, "Orchestration Pattern Comparison", 32, kicker="M5 · Choose control shape")
    table(s, 0.45, 1.08, 9.1, 4.05,
          ["Pattern", "Control owner", "Best fit", "Recovery / tradeoff"],
          [
              ("Sequential", "Workflow graph", "Ordered pipeline", "Simple checkpoints; serial latency"),
              ("Concurrent", "Fan-out / merge", "Independent evidence", "Partial results; merge semantics"),
              ("Handoff", "Current agent", "Triage / escalation", "Conversation transfer; route loops"),
              ("Group chat", "Manager / policy", "Debate / collaboration", "High token cost; convergence"),
              ("Magentic", "Dynamic manager", "Open-ended delegation", "Flexible; least predictable"),
          ], widths=[1.25, 1.75, 2.35, 3.75], font_size=9.1)
    textbox(s, 0.7, 5.18, 8.6, 0.2, "Rule: prefer the least dynamic pattern that satisfies the task and recovery requirements.", size=9.5, color=INDIGO, bold=True, align=PP_ALIGN.CENTER)

    # 33
    s = new_slide(prs); built.append(s); title(s, "Sequential Workflows", 33, kicker="M5 · Explicit order")
    code(s, 0.55, 1.13, 4.9, 3.8, '''workflow = (
    WorkflowBuilder(
        start_executor=intake,
        output_from="all",
    )
    .add_edge(intake, researcher)
    .add_edge(researcher, reviewer)
    .build()
)

result = await workflow.run(ticket)''', caption="agents as workflow nodes", size=8.9)
    flow(s, 1.4, ["INTAKE", "RESEARCH", "REVIEW"], x=5.78, total_w=3.35, height=0.62)
    card(s, 5.55, 2.35, 3.85, 1.0, "WHY", "Order is a correctness property; each stage narrows the next contract.", accent=INDIGO)
    card(s, 5.55, 3.58, 3.85, 1.0, "CHECKPOINT", "Persist stage input/output so restart resumes at the failed edge.", accent=GREEN, fill=PALE_GREEN)
    callout(s, 5.58, 4.78, 3.79, "Cost: total latency is additive.", accent=MAGENTA, size=9.2)

    # 34
    s = new_slide(prs); built.append(s); title(s, "Concurrent Fan-Out / Merge", 34, kicker="M5 · Parallel evidence")
    label(s, 0.7, 2.35, 1.4, "REQUEST", fill=INDIGO)
    for yy, name, accent in [(1.25, "SECURITY", RED), (2.35, "COST", CYAN), (3.45, "RELIABILITY", PURPLE)]:
        label(s, 3.4, yy, 1.55, name, fill=accent)
        arrow(s, 2.15, 2.52, 3.32, yy + 0.17, color=accent)
        arrow(s, 5.02, yy + 0.17, 6.38, 2.52, color=accent)
    label(s, 6.45, 2.35, 1.35, "MERGE", fill=GREEN)
    arrow(s, 7.85, 2.52, 9.0, 2.52, color=GREEN)
    label(s, 8.35, 3.05, 1.1, "DECIDE", fill=MAGENTA)
    card(s, 0.65, 4.25, 2.65, 0.9, "MERGE CONTRACT", "Keyed outputs; explicit missing branch state.", accent=INDIGO, body_size=9.4)
    card(s, 3.68, 4.25, 2.65, 0.9, "FAILURE POLICY", "Fail-fast, quorum, or degraded response.", accent=RED, fill=PALE_PINK, body_size=9.4)
    card(s, 6.7, 4.25, 2.65, 0.9, "RESOURCE POLICY", "Concurrency cap, timeout, cancellation.", accent=CYAN, fill=PALE_BLUE, body_size=9.4)

    # 35
    s = new_slide(prs); built.append(s); title(s, "Handoff, Group Chat, and Manager / Delegate Patterns", 35, kicker="M5 · Dynamic collaboration")
    table(s, 0.55, 1.15, 8.9, 3.7,
          ["Mode", "Control transfer", "Required invariant", "Watch for"],
          [
              ("Handoff", "One specialist becomes active", "Full conversation travels", "Ping-pong routes; lost ownership"),
              ("Group chat", "Shared conversation under manager", "Turn-selection + stop rule", "Echo chambers; token explosion"),
              ("Manager/delegate", "Manager plans and assigns work", "Task IDs + result contracts", "Manager bottleneck; weak verification"),
              ("Magentic", "Manager replans dynamically", "Budgets + durable progress", "Non-reproducible trajectories"),
          ], widths=[1.35, 2.2, 2.55, 2.8], font_size=9.1)
    callout(s, 0.75, 4.98, 8.5, "Handoff participants require per-service-call history persistence so control transfer carries complete state.", accent=INDIGO, size=9.5)

    # 36
    s = new_slide(prs); built.append(s); title(s, "Evaluation Loop and Golden Datasets", 36, kicker="M6 · Improve on purpose")
    flow(s, 1.16, ["DEFINE CHECKS", "RUN", "INSPECT FAILURES", "IMPROVE", "RERUN"])
    card(s, 0.65, 2.35, 4.1, 2.25, "GOLDEN CASE", "Query\nExpected output / fact groups\nMinimum required tool calls\nForbidden claims\nFrozen time and fixtures\nKnown refusal / escalation cases", accent=INDIGO, body_size=10.8)
    card(s, 5.25, 2.35, 4.1, 2.25, "DATASET HYGIENE", "Version beside code\nReview diffs in PRs\nDo not over-specify trajectories\nFail loudly on unknown cases\nReuse captured responses\nTrack per-slice regressions", accent=CYAN, fill=PALE_BLUE, body_size=10.8)
    callout(s, 0.75, 4.85, 8.5, "Freeze volatile facts: a test set whose right answer changes overnight is not golden.", accent=MAGENTA, fill=PALE_PINK)

    # 37
    s = new_slide(prs); built.append(s); title(s, "Deterministic Evaluators vs LLM Judges", 37, kicker="M6 · CI policy")
    table(s, 0.55, 1.12, 5.15, 3.8,
          ["Dimension", "Deterministic", "LLM judge"],
          [
              ("Checks", "Facts, schema, tools", "Groundedness, quality"),
              ("Repeatability", "High", "Model/version drift"),
              ("Debuggability", "Exact failure", "Rubric + rationale"),
              ("Cost", "Low", "Extra model calls"),
              ("CI role", "Gate build", "Inform analysis"),
          ], widths=[1.45, 1.85, 1.85], font_size=8.9)
    code(s, 5.95, 1.12, 3.45, 2.65, '''expected = [
    ExpectedToolCall(
        "get_order",
        {"order_id": "CR-1042"},
    )
]

result.raise_for_status()  # CI gate''', caption="minimum required behavior", size=8.25)
    callout(s, 5.98, 4.03, 3.39, "Judge off-format? Fail closed in the report—not the build.", accent=MAGENTA, fill=PALE_PINK, size=8.7)
    callout(s, 0.75, 5.05, 8.5, "Policy: deterministic checks gate CI; LLM judges expose nuanced quality trends.", accent=INDIGO)

    # 38
    s = new_slide(prs); built.append(s); title(s, "Production Guardrails and Operational Readiness", 38, kicker="M7 · Ship safely")
    table(s, 0.55, 1.12, 8.9, 3.82,
          ["Layer", "Guardrail", "Operational signal"],
          [
              ("Input", "Size, tenant auth, injection screening", "Reject rate / category"),
              ("Loop", "Iteration, time, token, and cost budgets", "Stop reason / budget burn"),
              ("Tools", "Allowlist, approval, timeout, idempotency", "Latency / denied calls / retries"),
              ("State", "Encryption, retention, tenant-scoped restore", "Persist failures / restore age"),
              ("Output", "Grounding, redaction, policy validation", "Blocked claims / citations"),
              ("Platform", "Circuit breakers, queues, rollout + rollback", "SLOs / saturation / release"),
          ], widths=[1.1, 4.35, 3.45], font_size=9.0)
    textbox(s, 0.7, 5.1, 8.6, 0.22, "Failure handling is product behavior: partial answer, retry, compensate, pause, or escalate—choose explicitly.", size=9.2, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 39
    s = new_slide(prs); built.append(s); title(s, "End-to-End Reference Architecture / Capstone Assembly", 39, kicker="M8 · Reference system")
    # Main architecture, all native shapes.
    label(s, 0.45, 1.45, 1.25, "CHANNELS", fill=PURPLE)
    textbox(s, 0.45, 1.88, 1.25, 1.0, "Web / API\nA2A client\nBatch / event", size=9.2, color=GRAY, align=PP_ALIGN.CENTER)
    label(s, 2.05, 1.45, 1.35, "HOST", fill=INDIGO)
    textbox(s, 2.0, 1.88, 1.45, 1.0, "Auth + tenancy\nQueue / cancel\nStream output", size=9.2, color=GRAY, align=PP_ALIGN.CENTER)
    box(s, 3.8, 1.25, 3.1, 2.45, fill=PALE, line=MAGENTA)
    textbox(s, 4.0, 1.45, 2.7, 0.35, "HARNESS AGENT", size=14, color=MAGENTA, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, 4.0, 1.95, 2.7, 1.4, "Todo + mode\nContext + compaction\nTools + approvals\nSkills + memory\nBackground agents", size=10, color=GRAY, align=PP_ALIGN.CENTER)
    label(s, 7.3, 1.1, 1.45, "MODEL", fill=CYAN)
    label(s, 7.3, 2.0, 1.45, "TOOLS / MCP", fill=GREEN)
    label(s, 7.3, 2.9, 1.45, "WORKFLOWS", fill=PURPLE)
    arrow(s, 1.72, 2.05, 2.0, 2.05, color=PURPLE)
    arrow(s, 3.47, 2.05, 3.75, 2.05, color=INDIGO)
    arrow(s, 6.95, 1.55, 7.25, 1.28, color=CYAN)
    arrow(s, 6.95, 2.25, 7.25, 2.17, color=GREEN)
    arrow(s, 6.95, 2.95, 7.25, 3.07, color=PURPLE)
    label(s, 0.65, 4.3, 1.6, "DURABLE STATE", fill=INDIGO)
    textbox(s, 2.4, 4.32, 2.0, 0.3, "history · files · checkpoints", size=9, color=GRAY)
    label(s, 4.7, 4.3, 1.45, "OTEL", fill=MAGENTA)
    textbox(s, 6.3, 4.32, 2.7, 0.3, "traces · usage · eval evidence · alerts", size=9, color=GRAY)
    arrow(s, 5.25, 3.75, 5.25, 4.23, color=MAGENTA)
    callout(s, 0.75, 4.92, 8.5, "Deploy as A2A, Azure Functions, or container—the state and provider boundaries remain explicit.", accent=CYAN, fill=PALE_BLUE, size=9.5)

    # 40
    s = new_slide(prs); built.append(s); title(s, "Key Takeaways and Next Steps", 40, kicker="Close")
    cards = [
        ("1 · LOOP", "Agent behavior emerges across calls. Instrument and bound the loop.", INDIGO, PALE),
        ("2 · STATE", "Make local history authoritative and persist after every service call.", MAGENTA, PALE_PINK),
        ("3 · CONTEXT", "Use progressive disclosure for skills and memory; compact deliberately.", CYAN, PALE_BLUE),
        ("4 · CONTROL", "Use workflows when order, recovery, or typed transitions are correctness properties.", PURPLE, PALE),
        ("5 · EVIDENCE", "Gate CI deterministically; use judges for analysis and trend detection.", GREEN, PALE_GREEN),
        ("6 · OPERATE", "Ship budgets, approvals, telemetry, failure policy, and hosting together.", RED, PALE_PINK),
    ]
    for i, (heading, body, accent, fill) in enumerate(cards):
        x = 0.6 + (i % 3) * 3.0
        y = 1.2 + (i // 3) * 1.65
        card(s, x, y, 2.7, 1.35, heading, body, accent=accent, fill=fill,
             heading_size=11, body_size=9.5)
    callout(s, 0.75, 4.72, 8.5, "Next: build M1 → add tools → engineer context → assemble the harness → evaluate → host.", accent=INDIGO)
    textbox(s, 0.75, 5.27, 8.5, 0.18, "monuminu.github.io/DHS-workshop26", size=9, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

    if len(built) != 40:
        raise RuntimeError(f"Expected 40 slides, built {len(built)}")

    # Remove the old reference slides after all clones and relationships exist.
    for _ in range(template_count):
        slide_id = prs.slides._sldIdLst[0]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[0]

    prs.core_properties.title = "Building Agents with the Microsoft Agent Framework"
    prs.core_properties.subject = "DHS 2026 deep technical workshop"
    prs.core_properties.author = "Manoranjan Rajguru"
    prs.core_properties.comments = "Generated from the DHS-workshop26 workshop content."
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
